# -*- coding: utf-8 -*-
"""
无 OCR 版：批量把 .doc/.docx/.ppt/.pptx/.pdf 中的可编辑文字提取成 txt
输入 / 输出路径直接写在代码里，双击即可运行
"""
import os
import sys
import logging
from pathlib import Path
import win32com.client as win32
from docx import Document
from pptx import Presentation
import PyPDF2

# =============== ① 在这里改你自己的文件夹 ===============
INPUT_FOLDER = r"D:\研究生\项目组\科研\小模型+RAG暑期实验\数据集\3054个word教案"  # 待扫描根目录
OUTPUT_FOLDER = r"D:\研究生\项目组\科研\小模型+RAG暑期实验\数据集\3054个txt教案"  # 结果保存根目录
# =========================================================

logging.basicConfig(level=logging.INFO, format="%(asctime)s | %(levelname)s | %(message)s")


# ---------------- 工具函数 ----------------
def safe_com_close(app, doc=None):
    try:
        if doc:
            doc.Close(False)
        if app:
            app.Quit()
    except Exception as e:
        logging.warning("关闭 COM 对象异常：%s", e)


def remove_non_printable(text: str) -> str:
    return "".join(c for c in text if c.isprintable())


# ---------------- 读取函数 ----------------
def read_doc(fp: str) -> str:
    text = ""
    word = None
    doc = None
    try:
        word = win32.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(fp, Encoding="utf-8")
        for para in doc.Paragraphs:
            text += para.Range.Text.strip().replace("\n", "")
        return remove_non_printable(text)
    except Exception as e:
        logging.error("读 .doc 失败：%s  %s", fp, e)
        return ""
    finally:
        safe_com_close(word, doc)
        del word


def read_docx(fp: str) -> str:
    doc = Document(fp)
    full_text = []
    for p in doc.paragraphs:
        full_text.append("".join(p.text.split()))
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                full_text.append(cell.text)
    return "".join(full_text)


def read_ppt(fp: str) -> str:
    text = ""
    powerpoint = None
    pres = None
    try:
        powerpoint = win32.gencache.EnsureDispatch("PowerPoint.Application")
        pres = powerpoint.Presentations.Open(fp, WithWindow=False)
        for slide in pres.Slides:
            for shp in slide.Shapes:
                if shp.HasTextFrame:
                    text += shp.TextFrame.TextRange.Text
        return text
    except Exception as e:
        logging.error("读 .ppt 失败：%s  %s", fp, e)
        return ""
    finally:
        safe_com_close(powerpoint, pres)
        del powerpoint


def read_pptx(fp: str) -> str:
    prs = Presentation(fp)
    text = ""
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame:
                for p in shp.text_frame.paragraphs:
                    for run in p.runs:
                        text += run.text
    return text


def read_pdf(fp: str) -> str:
    text = ""
    try:
        with open(fp, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
    except Exception as e:
        logging.error("读 .pdf 失败：%s  %s", fp, e)
    return text


# ---------------- 主逻辑 ----------------
def extract_file(fp: Path, out_dir: Path):
    suffix = fp.suffix.lower()
    body = ""
    try:
        if suffix == ".doc":
            body = read_doc(str(fp))
        elif suffix == ".docx":
            body = read_docx(str(fp))
        elif suffix == ".ppt":
            body = read_ppt(str(fp))
        elif suffix == ".pptx":
            body = read_pptx(str(fp))
        elif suffix == ".pdf":
            body = read_pdf(str(fp))
        else:
            return
    except Exception as e:
        logging.error("提取异常：%s  %s", fp, e)
        return

    # 镜像目录结构
    rel_path = fp.relative_to(INPUT_ROOT)
    out_sub = out_dir / rel_path.parent
    out_sub.mkdir(parents=True, exist_ok=True)

    txt_file = out_sub / f"{fp.stem}.txt"
    try:
        txt_file.write_text(body, encoding="utf-8")
        logging.info("已写入：%s", txt_file)
    except Exception as e:
        logging.error("写 txt 失败：%s  %s", txt_file, e)


def walk_and_extract(in_root: Path, out_root: Path):
    for fp in in_root.rglob("*"):
        if fp.is_file() and fp.suffix.lower() in {".doc", ".docx", ".ppt", ".pptx", ".pdf"}:
            extract_file(fp, out_root)


# ---------------- 入口 ----------------
if __name__ == "__main__":
    INPUT_ROOT = Path(INPUT_FOLDER).resolve()
    OUTPUT_ROOT = Path(OUTPUT_FOLDER).resolve()

    if not INPUT_ROOT.is_dir():
        logging.error("输入目录不存在：%s", INPUT_ROOT)
        sys.exit(1)

    OUTPUT_ROOT.mkdir(parents=True, exist_ok=True)

    walk_and_extract(INPUT_ROOT, OUTPUT_ROOT)
    logging.info("全部处理完成！结果保存在：%s", OUTPUT_ROOT)