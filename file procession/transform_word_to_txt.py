import os
import re
from docx import Document
from pptx import Presentation
from pdfminer.high_level import extract_text
from olefile import OleFileIO


def docx_to_text(docx_path):
    """提取DOCX文本（保留段落）"""
    doc = Document(docx_path)
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def doc_to_text(doc_path):
    """提取DOC文本（旧版Word）"""
    try:
        ole = OleFileIO(doc_path)
        if ole.exists('WordDocument'):
            with ole.openstream('WordDocument') as f:
                text = f.read().decode('latin-1')
                return re.sub(r'[\x00-\x1F\x7F]+', ' ', text)  # 去除控制字符
        return ""
    except Exception as e:
        print(f"[警告] 解析{doc_path}失败: {str(e)}")
        return ""


def pptx_to_text(pptx_path):
    """提取PPTX文本"""
    prs = Presentation(pptx_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def ppt_to_text(ppt_path):
    """提取PPT文本（旧版PPT）"""
    # 旧版PPT解析较复杂，这里返回空字符串（实际项目可用comtypes调用PowerPoint COM接口）
    print(f"[提示] 旧版PPT({ppt_path})需手动转换或使用专业工具")
    return ""


def pdf_to_text(pdf_path):
    """提取PDF文本"""
    try:
        return extract_text(pdf_path)
    except Exception as e:
        print(f"[警告] 解析{pdf_path}失败: {str(e)}")
        return ""


def convert_file(input_path, output_folder):
    """根据扩展名选择转换器"""
    ext = os.path.splitext(input_path)[1].lower()
    filename = os.path.basename(input_path)

    try:
        if ext == '.docx':
            text = docx_to_text(input_path)
        elif ext == '.doc':
            text = doc_to_text(input_path)
        elif ext == '.pptx':
            text = pptx_to_text(input_path)
        elif ext == '.ppt':
            text = ppt_to_text(input_path)
        elif ext == '.pdf':
            text = pdf_to_text(input_path)
        else:
            print(f"[跳过] 不支持的文件类型: {filename}")
            return False

        output_path = os.path.join(output_folder, f"{os.path.splitext(filename)[0]}.txt")
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        print(f"[成功] {filename} → {os.path.basename(output_path)}")
        return True
    except Exception as e:
        print(f"[失败] {filename}: {str(e)}")
        return False


def batch_convert(input_folder, output_folder):
    """批量转换入口函数"""
    if not os.path.exists(input_folder):
        raise FileNotFoundError(f"输入文件夹不存在: {input_folder}")

    os.makedirs(output_folder, exist_ok=True)
    success = 0

    for filename in sorted(os.listdir(input_folder)):
        input_path = os.path.join(input_folder, filename)
        if os.path.isfile(input_path):
            if convert_file(input_path, output_folder):
                success += 1

    print(f"\n转换完成！成功: {success} 个文件 | 输出目录: {os.path.abspath(output_folder)}")


if __name__ == "__main__":
    # 配置路径
    input_dir =r"D:\研究生\项目组\科研\小模型+RAG暑期实验\数据集\几万条word教案" # 存放所有文档的文件夹
    output_dir =r"D:\研究生\项目组\科研\小模型+RAG暑期实验\数据集\txt数据集"   # 输出文件夹

    batch_convert(input_dir, output_dir)








